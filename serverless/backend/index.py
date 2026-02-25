"""MBA Copilot - FastAPI Backend.

A RAG-powered document Q&A system for MBA students.
"""

from __future__ import annotations

import csv
import io
import os
import random
import re
import string
import time
from datetime import datetime, timezone
from typing import TYPE_CHECKING, Annotated, Any, cast

import fitz  # PyMuPDF
import tiktoken
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pydantic import BaseModel

if TYPE_CHECKING:
    from openai import OpenAI
    from openai.types.chat import ChatCompletionMessageParam
    from pinecone import Index

# =============================================================================
# App
# =============================================================================
load_dotenv()
app = FastAPI(title="MBA Copilot API", root_path="/backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # tighten in prod
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =============================================================================
# Configuration
# =============================================================================


class Config:
    """Configuration settings for the MBA Copilot application."""

    # OpenAI
    OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
    OPENAI_BASE_URL = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")
    EMBEDDING_MODEL = "text-embedding-3-large"
    CHAT_MODEL = "gpt-4o-mini"
    EMBEDDING_DIMENSIONS = 1024

    # Pinecone
    PINECONE_API_KEY = os.environ.get("PINECONE_API_KEY")
    PINECONE_INDEX = os.environ.get("PINECONE_INDEX", "mba-copilot")

    # RAG Settings (token-based)
    # Larger chunks to keep more context together
    CHUNK_TOKENS_DOCS = 800
    CHUNK_OVERLAP_TOKENS_DOCS = 150

    # Retrieval settings
    RETRIEVAL_TOP_K = 20  # Retrieve more candidates
    CONTEXT_MAX_CHUNKS = 8  # Pass more context to LLM
    MIN_SCORE = 0.25  # Lower threshold to be more inclusive


config = Config()

# =============================================================================
# Clients (lazy init)
# =============================================================================

_openai_client: OpenAI | None = None
_pinecone_index: Index | None = None


def get_openai() -> OpenAI:
    """Get or initialize the OpenAI client."""
    global _openai_client
    if _openai_client is None:
        if not config.OPENAI_API_KEY:
            raise RuntimeError("OPENAI_API_KEY is not set")

        from openai import OpenAI

        _openai_client = OpenAI(
            api_key=config.OPENAI_API_KEY,
            base_url=config.OPENAI_BASE_URL,
        )

    return _openai_client


def get_pinecone_index() -> Index:
    """Get or initialize the Pinecone index."""
    global _pinecone_index
    if _pinecone_index is None:
        if not config.PINECONE_API_KEY:
            raise RuntimeError("PINECONE_API_KEY is not set")

        from pinecone import Pinecone

        pc = Pinecone(api_key=config.PINECONE_API_KEY)
        _pinecone_index = pc.Index(config.PINECONE_INDEX)

    return _pinecone_index


# =============================================================================
# Token Utilities
# =============================================================================


def num_tokens(text: str, model: str | None = None) -> int:
    """Count tokens in text using tiktoken."""
    encoding_model = model or config.EMBEDDING_MODEL
    try:
        enc = tiktoken.encoding_for_model(encoding_model)
    except KeyError:
        # Fallback to cl100k_base for unknown models
        enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))


def chunk_by_tokens(
    text: str,
    chunk_tokens: int,
    overlap_tokens: int,
    model: str | None = None,
) -> list[str]:
    """Split text into chunks by token count (not characters)."""
    text = text.replace("\r\n", "\n").strip()
    if not text:
        return []

    if overlap_tokens >= chunk_tokens:
        raise ValueError("overlap_tokens must be < chunk_tokens")

    encoding_model = model or config.EMBEDDING_MODEL
    try:
        enc = tiktoken.encoding_for_model(encoding_model)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")

    tokens = enc.encode(text)

    # If text fits in one chunk, return as-is
    if len(tokens) <= chunk_tokens:
        return [text]

    chunks: list[str] = []
    start = 0

    while start < len(tokens):
        end = min(start + chunk_tokens, len(tokens))
        chunk_text = enc.decode(tokens[start:end]).strip()
        if chunk_text:
            chunks.append(chunk_text)

        if end >= len(tokens):
            break

        # Move start forward, accounting for overlap
        start = max(0, end - overlap_tokens)

    return chunks


# =============================================================================
# Document Processing
# =============================================================================


def _extract_pdf_text_best_fidelity(content: bytes) -> str:
    """Extract text from PDF with best fidelity using block coordinates.

    - get_text("blocks") -> includes bounding boxes
    - sort blocks by y then x with a small y tolerance to keep lines aligned
    - preserve page boundaries
    """
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        pages_out: list[str] = []
        y_tol = 3.0  # points

        for page in doc:
            blocks: Any = page.get_text("blocks")
            clean_blocks: list[Any] = []

            for b in blocks:
                if (
                    isinstance(b, (tuple, list))
                    and len(b) >= 5
                    and isinstance(b[4], str)
                    and b[4].strip()
                ):
                    clean_blocks.append(b)

            clean_blocks.sort(key=lambda b: (round(float(b[1]) / y_tol), float(b[0])))

            page_text = "\n".join(str(b[4]).rstrip() for b in clean_blocks).strip()
            if page_text:
                pages_out.append(page_text)

        return "\n\n".join(pages_out).strip()
    finally:
        doc.close()


def _extract_docx_text(content: bytes) -> str:
    """Extract text from DOCX file."""
    doc = Document(io.BytesIO(content))
    parts: list[str] = []

    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)

    for table in doc.tables:
        for row in table.rows:
            line = "\t".join(cell.text.strip() for cell in row.cells).rstrip()
            if line.strip():
                parts.append(line)

    return "\n".join(parts).strip()


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX file."""
    prs = Presentation(io.BytesIO(content))
    slides_out: list[str] = []

    for si, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"--- Slide {si} ---"]

        for shape in slide.shapes:
            # python-pptx is dynamic; stubs are conservative.
            s = cast(Any, shape)

            if hasattr(s, "text_frame") and s.text_frame:
                txt = (s.text_frame.text or "").strip()
                if txt:
                    parts.append(txt)

            if hasattr(s, "table") and s.table:
                for row in s.table.rows:
                    line = "\t".join(cell.text.strip() for cell in row.cells).rstrip()
                    if line.strip():
                        parts.append(line)

        # Speaker notes
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_txt = (notes_slide.notes_text_frame.text or "").strip()
                if notes_txt:
                    parts.append("[Notes]\n" + notes_txt)
        except Exception:
            pass

        slides_out.append("\n".join(parts).strip())

    return "\n\n".join(s for s in slides_out if s).strip()


def _extract_csv_structured(content: bytes) -> list[dict[str, Any]]:
    """Extract CSV as row-based chunks with column headers.

    Each row becomes a chunk formatted as: "ColA: valA | ColB: valB | ..."
    """
    text = content.decode("utf-8-sig", errors="replace")
    reader = csv.DictReader(io.StringIO(text))

    rows: list[dict[str, Any]] = []
    for idx, row in enumerate(reader, start=1):
        # Format: "ColA: valA | ColB: valB"
        chunk_text = " | ".join(f"{k}: {v}" for k, v in row.items() if v and v.strip())
        if chunk_text.strip():
            rows.append({
                "row_number": idx,
                "text": chunk_text,
            })

    return rows


def _extract_pdf_with_pages(content: bytes) -> list[dict[str, Any]]:
    """Extract PDF with page-level metadata.

    Returns list of dicts with: page_number, text
    """
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        pages: list[dict[str, Any]] = []
        y_tol = 3.0  # points

        for page_num, page in enumerate(doc, start=1):
            blocks: Any = page.get_text("blocks")
            clean_blocks: list[Any] = []

            for b in blocks:
                if (
                    isinstance(b, (tuple, list))
                    and len(b) >= 5
                    and isinstance(b[4], str)
                    and b[4].strip()
                ):
                    clean_blocks.append(b)

            clean_blocks.sort(key=lambda b: (round(float(b[1]) / y_tol), float(b[0])))

            page_text = "\n".join(str(b[4]).rstrip() for b in clean_blocks).strip()
            if page_text:
                pages.append({
                    "page_number": page_num,
                    "text": page_text,
                })

        return pages
    finally:
        doc.close()


def extract_structured_chunks(file: UploadFile) -> list[dict[str, Any]]:
    """Extract file into structured chunks with metadata.

    Simple token-based chunking for all file types.
    Returns list of dicts with 'text' and 'chunk_index'.
    """
    content = file.file.read()
    try:
        file.file.seek(0)
    except Exception:
        pass

    if not isinstance(file.filename, str) or not file.filename:
        raise ValueError("Uploaded file has no filename")

    filename = file.filename.lower()

    # Extract text based on file type
    if filename.endswith(".pptx"):
        text = _extract_pptx_text(content)
    elif filename.endswith(".csv"):
        # For CSV, just treat as plain text for now
        # TODO: Revisit row-based chunking with batching when we have more time
        #
        # # For CSV: Use row-based chunking with batching to balance context and performance
        # # Convert each row to "Col1: val1 | Col2: val2" format, then batch rows together
        # text_content = content.decode("utf-8-sig", errors="replace")
        # reader = csv.DictReader(io.StringIO(text_content))
        #
        # row_texts: list[str] = []
        # for row in reader:
        #     # Format: "ColA: valA | ColB: valB"
        #     row_text = " | ".join(f"{k}: {v}" for k, v in row.items() if v and v.strip())
        #     if row_text.strip():
        #         row_texts.append(row_text)
        #
        # if not row_texts:
        #     return []
        #
        # # Batch rows into chunks (target ~400-500 tokens per chunk for good retrieval)
        # # Average row is ~50-100 tokens, so batch 5-10 rows per chunk
        # chunks: list[str] = []
        # current_chunk: list[str] = []
        # current_tokens = 0
        # target_tokens = 450  # Sweet spot for retrieval
        #
        # for row_text in row_texts:
        #     row_tokens = num_tokens(row_text)
        #
        #     if current_tokens + row_tokens > target_tokens and current_chunk:
        #         # Chunk is full, save it and start new one
        #         chunks.append("\n".join(current_chunk))
        #         current_chunk = [row_text]
        #         current_tokens = row_tokens
        #     else:
        #         # Add to current chunk
        #         current_chunk.append(row_text)
        #         current_tokens += row_tokens
        #
        # # Don't forget the last chunk
        # if current_chunk:
        #     chunks.append("\n".join(current_chunk))
        #
        # return [{"text": chunk, "chunk_index": i} for i, chunk in enumerate(chunks)]

        text = content.decode("utf-8-sig", errors="replace")
    elif filename.endswith(".pdf"):
        text = _extract_pdf_text_best_fidelity(content)
    elif filename.endswith(".docx"):
        text = _extract_docx_text(content)
    elif filename.endswith((".txt", ".md")):
        text = content.decode("utf-8-sig", errors="replace")
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {filename}")

    if not text.strip():
        return []

    # Chunk everything with token-based chunking
    text_chunks = chunk_by_tokens(
        text,
        chunk_tokens=config.CHUNK_TOKENS_DOCS,
        overlap_tokens=config.CHUNK_OVERLAP_TOKENS_DOCS,
    )

    return [{"text": chunk, "chunk_index": i} for i, chunk in enumerate(text_chunks)]


def generate_document_id() -> str:
    """Generate a unique document ID."""
    return f"doc_{int(time.time())}_{''.join(random.choices(string.ascii_lowercase, k=6))}"


# =============================================================================
# Embeddings
# =============================================================================


def generate_embedding(text: str) -> list[float]:
    """Generate embedding for a single text string."""
    client = get_openai()
    response = client.embeddings.create(
        model=config.EMBEDDING_MODEL, input=text, dimensions=config.EMBEDDING_DIMENSIONS
    )
    return response.data[0].embedding


async def generate_embeddings_batch(texts: list[str]) -> list[list[float]]:
    """Generate embeddings for multiple texts using parallel individual requests.

    Note: CBS endpoint blocks batch requests via Cloudflare, so we send
    individual requests in parallel instead.

    TODO: Revert to batch API when CBS IT enables batch embedding requests.
    Original implementation:
        client = get_openai()
        response = client.embeddings.create(
            model=config.EMBEDDING_MODEL, input=texts, dimensions=config.EMBEDDING_DIMENSIONS
        )
        return [d.embedding for d in response.data]
    """
    import asyncio

    from openai import AsyncOpenAI

    # Create async client with same config as sync client
    async_client = AsyncOpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL if config.OPENAI_BASE_URL else None,
    )

    async def get_single_embedding(text: str) -> list[float]:
        """Get embedding for a single text."""
        response = await async_client.embeddings.create(
            model=config.EMBEDDING_MODEL, input=text, dimensions=config.EMBEDDING_DIMENSIONS
        )
        return response.data[0].embedding

    # Run all requests in parallel
    embeddings = await asyncio.gather(*[get_single_embedding(text) for text in texts])
    return list(embeddings)


# =============================================================================
# Pinecone Operations
# =============================================================================


def store_chunks(chunks: list[dict[str, Any]]) -> None:
    """Store document chunks in Pinecone vector database."""
    index = get_pinecone_index()

    batch_size = 100
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i : i + batch_size]
        vectors = [
            {"id": c["id"], "values": c["embedding"], "metadata": c["metadata"]} for c in batch
        ]
        index.upsert(vectors=vectors)


def query_similar(
    embedding: list[float],
    top_k: int | None = None,
    document_ids: list[str] | None = None,
) -> list[dict[str, Any]]:
    """Query Pinecone for similar document chunks."""
    index = get_pinecone_index()

    query_filter = None
    if document_ids:
        query_filter = {"document_id": {"$in": document_ids}}

    results = index.query(
        vector=embedding,
        top_k=top_k or config.RETRIEVAL_TOP_K,
        include_metadata=True,
        filter=query_filter,
    )

    return [
        {
            "id": m.id,
            "score": m.score,
            "text": (m.metadata or {}).get("text", ""),
            "filename": (m.metadata or {}).get("filename", ""),
            "document_id": (m.metadata or {}).get("document_id", ""),
            "metadata": m.metadata,
        }
        for m in results.matches
    ]


def delete_document(document_id: str) -> None:
    """Delete all chunks for a document from Pinecone."""
    index = get_pinecone_index()
    index.delete(filter={"document_id": {"$eq": document_id}})


def list_documents() -> list[dict[str, Any]]:
    """Best-effort listing using the 'is_first_chunk' marker.

    NOTE: This depends on Pinecone supporting metadata filtering.
    """
    index = get_pinecone_index()

    # Reduce top_k to avoid timeouts - 100 documents should be enough
    # for most use cases and is much faster than 1000
    results = index.query(
        vector=[0.0] * config.EMBEDDING_DIMENSIONS,
        top_k=100,
        include_metadata=True,
        filter={"is_first_chunk": {"$eq": True}},
    )

    documents: list[dict[str, Any]] = []
    for m in results.matches:
        md = m.metadata or {}
        documents.append(
            {
                "id": md.get("document_id"),
                "filename": md.get("filename"),
                "chunks": md.get("total_chunks", 1),
                "uploaded_at": md.get("uploaded_at", ""),
            }
        )

    return documents


# =============================================================================
# RAG Pipeline
# =============================================================================


def generate_answer(
    question: str,
    context: str,
    history: list[dict[str, Any]] | None = None,
    chat_model: str | None = None,
    system_prompt: str | None = None,
) -> str:
    """Generate an answer using OpenAI's chat completion API."""
    client = get_openai()

    prompt = system_prompt or "You are a helpful AI assistant."
    model = chat_model or config.CHAT_MODEL

    # Build messages list with proper typing
    messages: list[dict[str, str]] = [{"role": "system", "content": prompt}]

    if context:
        messages.append(
            {
                "role": "system",
                "content": (
                    "Here is relevant information from the student's documents:\n\n"
                    f"{context}\n\n"
                    "Use this to answer the question. Cite sources when appropriate."
                ),
            }
        )
    else:
        messages.append(
            {
                "role": "system",
                "content": (
                    "No relevant documents were found. If needed, let the student know "
                    "they should upload relevant materials, but still try to help with general knowledge."
                ),
            }
        )

    # Include full history - OpenAI API handles token limits gracefully
    # by truncating from the beginning if needed
    if history:
        for msg in history:
            role = str(msg["role"])
            content = str(msg["content"])
            if role in ("user", "assistant", "system"):
                messages.append({"role": role, "content": content})

    messages.append({"role": "user", "content": question})

    # Cast to the proper type for OpenAI API
    response = client.chat.completions.create(
        model=model,
        messages=cast("list[ChatCompletionMessageParam]", messages),
        temperature=0.7,
        max_tokens=1000,
    )
    return response.choices[0].message.content or ""


# =============================================================================
# API Models
# =============================================================================


class ChatSettings(BaseModel):
    """Settings for chat completion and RAG retrieval.

    Note: top_k is now just for backwards compatibility.
    The system retrieves config.RETRIEVAL_TOP_K candidates and passes
    the best config.CONTEXT_MAX_CHUNKS to the LLM.
    """

    chat_model: str = "gpt-4o-mini"
    top_k: int = 15  # Kept for backwards compatibility
    min_score: float = 0.3
    system_prompt: str = "You are a helpful AI assistant."


class ChatRequest(BaseModel):
    """Request model for chat endpoint."""

    message: str
    history: list[dict[str, Any]] | None = None
    settings: ChatSettings | None = None
    document_ids: list[str] | None = None


class ChatResponse(BaseModel):
    """Response model for chat endpoint."""

    answer: str
    sources: list[dict[str, Any]]


# =============================================================================
# API Endpoints
# =============================================================================


@app.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest) -> ChatResponse:
    """Chat endpoint - answer questions using RAG retrieval."""
    try:
        settings = request.settings or ChatSettings()

        query_embedding = generate_embedding(request.message)

        # Retrieve more candidates than we'll use
        similar = query_similar(
            query_embedding,
            top_k=config.RETRIEVAL_TOP_K,
            document_ids=request.document_ids,
        )

        # Filter by minimum score
        relevant = [c for c in similar if float(c.get("score", 0.0)) >= settings.min_score]

        # If we have results, limit to best N for context
        if relevant:
            context_chunks = relevant[: config.CONTEXT_MAX_CHUNKS]
        elif similar:
            # Fallback: if min_score filtered everything, use top results anyway
            context_chunks = similar[: max(3, config.CONTEXT_MAX_CHUNKS // 2)]
        else:
            context_chunks = []

        # Build context
        if context_chunks:
            context = "\n\n---\n\n".join(
                [f"[Source: {c['filename']}]\n{c['text']}" for c in context_chunks]
            )
        else:
            context = ""

        answer = generate_answer(
            request.message,
            context,
            request.history,
            chat_model=settings.chat_model,
            system_prompt=settings.system_prompt,
        )

        # Return sources from context_chunks (what was actually used)
        sources = [
            {
                "text": c["text"],
                "score": c["score"],
                "filename": c["filename"],
                "document_id": c["document_id"],
                "metadata": c.get("metadata", {}),
            }
            for c in context_chunks
        ]

        return ChatResponse(answer=answer, sources=sources)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


class _FileObj:
    """Lightweight file-like object compatible with extract_structured_chunks."""

    def __init__(self, content: bytes, filename: str) -> None:
        self.file = io.BytesIO(content)
        self.filename = filename
        self.content_type = "application/octet-stream"


def _make_file_obj(content: bytes, filename: str) -> _FileObj:
    return _FileObj(content, filename)


async def _process_file(file_obj: Any, display_filename: str) -> dict[str, Any]:
    """Shared file processing: extract chunks, generate embeddings, store in Pinecone."""
    if not config.OPENAI_API_KEY:
        raise HTTPException(
            status_code=500,
            detail="OPENAI_API_KEY not configured. Please set environment variables in Vercel dashboard.",
        )
    if not config.PINECONE_API_KEY:
        raise HTTPException(
            status_code=500,
            detail="PINECONE_API_KEY not configured. Please set environment variables in Vercel dashboard.",
        )

    structured_chunks = extract_structured_chunks(file_obj)
    if not structured_chunks:
        raise HTTPException(status_code=400, detail="No content to process")

    # Check if document with same filename already exists and delete it
    existing_docs = list_documents()
    for doc in existing_docs:
        if doc.get("filename") == display_filename:
            print(f"Deleting existing document with filename: {display_filename}")
            delete_document(doc["id"])

    chunk_texts = [chunk["text"] for chunk in structured_chunks]
    embeddings = await generate_embeddings_batch(chunk_texts)

    document_id = generate_document_id()
    uploaded_at = datetime.now(timezone.utc).isoformat()

    chunks: list[dict[str, Any]] = []
    for i, (structured_chunk, embedding) in enumerate(
        zip(structured_chunks, embeddings, strict=False)
    ):
        chunks.append({
            "id": f"{document_id}_chunk_{i}",
            "embedding": embedding,
            "metadata": {
                "text": structured_chunk["text"],
                "document_id": document_id,
                "filename": display_filename,
                "chunk_index": i,
                "total_chunks": len(structured_chunks),
                "uploaded_at": uploaded_at,
                "is_first_chunk": i == 0,
            },
        })

    store_chunks(chunks)

    return {
        "success": True,
        "document_id": document_id,
        "filename": display_filename,
        "chunks": len(structured_chunks),
    }


@app.post("/upload")
async def upload(
    file: Annotated[UploadFile, File()],
    filename: Annotated[str | None, Form()] = None,
) -> dict[str, Any]:
    """Upload and process a document file."""
    try:
        display_filename = filename or file.filename or "unknown"
        return await _process_file(file, display_filename)
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_detail = f"{str(e)}\n\nTraceback:\n{traceback.format_exc()}"
        raise HTTPException(status_code=500, detail=error_detail) from e


@app.get("/documents")
async def get_documents() -> dict[str, Any]:
    """Get list of all uploaded documents."""
    try:
        documents = list_documents()
        return {"documents": documents}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


@app.delete("/documents/{document_id}")
async def remove_document(document_id: str) -> dict[str, bool]:
    """Delete a document and all its chunks."""
    try:
        delete_document(document_id)
        return {"success": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


@app.post("/upload-from-url")
async def upload_from_url(request: dict[str, Any]) -> dict[str, Any]:
    """Download a file from a URL and process it."""
    try:
        url = request.get("url")
        filename = request.get("filename")

        if not url or not filename:
            raise HTTPException(status_code=400, detail="Missing url or filename")

        import httpx

        async with httpx.AsyncClient(timeout=300.0) as client:
            response = await client.get(url)
            response.raise_for_status()
            content = response.content

        fake_file = _make_file_obj(content, filename)
        return await _process_file(fake_file, filename)

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_detail = f"{str(e)}\n\nTraceback:\n{traceback.format_exc()}"
        raise HTTPException(status_code=500, detail=error_detail) from e


@app.post("/upload-from-urls")
async def upload_from_urls(request: dict[str, Any]) -> dict[str, Any]:
    """Download file parts from multiple blob URLs, concatenate, and process.

    Used by the chunked upload flow: each part was uploaded as an individual
    small blob. This endpoint downloads them all in parallel, concatenates
    in order, and processes the assembled file.
    """
    try:
        urls: list[str] = request.get("urls", [])
        filename: str | None = request.get("filename")

        if not urls or not filename:
            raise HTTPException(status_code=400, detail="Missing urls or filename")

        import asyncio

        import httpx

        async with httpx.AsyncClient(timeout=300.0) as client:
            responses = await asyncio.gather(
                *[client.get(url) for url in urls]
            )

        # Concatenate parts in order (urls are already sorted by the caller)
        content = b"".join(r.content for r in responses)
        print(f"[upload-from-urls] Downloaded {len(urls)} parts ({len(content) / 1024 / 1024:.2f} MB)")

        fake_file = _make_file_obj(content, filename)
        return await _process_file(fake_file, filename)

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_detail = f"{str(e)}\n\nTraceback:\n{traceback.format_exc()}"
        raise HTTPException(status_code=500, detail=error_detail) from e


@app.get("/health")
async def health() -> dict[str, str]:
    """Health check endpoint."""
    return {"status": "ok"}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
